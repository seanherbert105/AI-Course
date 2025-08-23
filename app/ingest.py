import os
import weaviate
from weaviate.util import generate_uuid5
from weaviate.collections.classes.config import DataType
import pdfplumber
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import textract

# ====== CONFIG ======
WEAVIATE_URL = "http://weaviate:8080"
DOCS_DIR = "/app/docs"
CLASS_NAME = "Eval"  # collection name in Weaviate
# ====================

# ✅ Connect with v4 client
client = weaviate.connect_to_custom(
    http_host="weaviate",
    http_port=8080,
    http_secure=False,
    grpc_host="weaviate",
    grpc_port=50051,
    grpc_secure=False,
)

print("Weaviate ready? ->", client.is_ready())

def extract_text_from_file(filepath: str) -> str | None:
    """Extract text from various file types."""
    ext = filepath.lower()
    try:
        if ext.endswith((".txt", ".md")):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        if ext.endswith(".pdf"):
            text = ""
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text

        if ext.endswith(".docx"):
            doc = DocxDocument(filepath)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext.endswith(".doc"):
            return textract.process(filepath).decode("utf-8", errors="ignore")

        if ext.endswith(".csv"):
            df = pd.read_csv(filepath)
            return df.to_string(index=False)

        if ext.endswith(".xlsx"):
            df = pd.read_excel(filepath)
            return df.to_string(index=False)

        if ext.endswith(".pptx"):
            prs = Presentation(filepath)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        if ext.endswith(".ppt"):
            return textract.process(filepath).decode("utf-8", errors="ignore")

    except Exception as e:
        print(f"❌ Error reading {filepath}: {e}")
        return None

    return None


def ensure_schema() -> None:
    # v4 returns list[str], not objects
    existing_collections = client.collections.list_all()

    if CLASS_NAME not in existing_collections:
        client.collections.create(
            name=CLASS_NAME,
            # You can omit vectorizer_config to use server defaults (e.g. text2vec-ollama)
            properties=[
                {"name": "title", "data_type": DataType.TEXT},
                {"name": "content", "data_type": DataType.TEXT},
            ],
        )
        print(f"✅ Created collection: {CLASS_NAME}")
    else:
        print(f"ℹ️ Collection already exists: {CLASS_NAME}")


def ingest_docs() -> None:
    ensure_schema()
    collection = client.collections.get(CLASS_NAME)

    for root, _, files in os.walk(DOCS_DIR):
        for fname in files:
            path = os.path.join(root, fname)
            text = extract_text_from_file(path)

            if not text or not text.strip():
                print(f"⚠️ Skipping unsupported or empty file: {fname}")
                continue

            # v4 insert
            collection.data.insert(
                properties={"title": fname, "content": text},
                uuid=generate_uuid5(fname + text),
            )
            print(f"📄 Ingested: {fname}")


if __name__ == "__main__":
    try:
        ingest_docs()
    finally:
        client.close()