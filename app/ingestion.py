import os
import weaviate
from weaviate.util import generate_uuid5
import pdfplumber
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import textract

# ====== CONFIG ======
WEAVIATE_URL = "http://weaviate:8080"
DOCS_DIR = "/app/docs"
CLASS_NAME = "Document"  # Schema class name in Weaviate
OLLAMA_MODEL = "mistral"  # Model running in Ollama
# ====================

def extract_text_from_file(filepath):
    """Extract text from various file types."""
    ext = filepath.lower()

    try:
        if ext.endswith((".txt", ".md")):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        elif ext.endswith(".pdf"):
            text = ""
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text

        elif ext.endswith(".docx"):
            doc = DocxDocument(filepath)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext.endswith(".doc"):
            return textract.process(filepath).decode("utf-8", errors="ignore")

        elif ext.endswith(".csv"):
            df = pd.read_csv(filepath)
            return df.to_string(index=False)

        elif ext.endswith(".xlsx"):
            df = pd.read_excel(filepath)
            return df.to_string(index=False)

        elif ext.endswith(".pptx"):
            prs = Presentation(filepath)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        elif ext.endswith(".ppt"):
            return textract.process(filepath).decode("utf-8", errors="ignore")

    except Exception as e:
        print(f"❌ Error reading {filepath}: {e}")
        return None

    return None

def ensure_schema(client):
    """Create a Weaviate class for documents if it doesn't exist."""
    schema = {
        "class": CLASS_NAME,
        "vectorizer": "text2vec-ollama",
        "moduleConfig": {
            "text2vec-ollama": {
                "model": OLLAMA_MODEL,
                "options": {}
            }
        },
        "properties": [
            {"name": "filename", "dataType": ["text"]},
            {"name": "content", "dataType": ["text"]}
        ]
    }

    existing_classes = [c["class"] for c in client.schema.get().get("classes", [])]
    if CLASS_NAME not in existing_classes:
        client.schema.create_class(schema)
        print(f"✅ Created class '{CLASS_NAME}' in Weaviate.")
    else:
        print(f"ℹ️ Class '{CLASS_NAME}' already exists.")

def ingest_docs():
    client = weaviate.Client(WEAVIATE_URL)
    ensure_schema(client)

    for root, _, files in os.walk(DOCS_DIR):
        for file in files:
            path = os.path.join(root, file)
            text = extract_text_from_file(path)

            if not text:
                print(f"⚠️ Skipping unsupported or unreadable file: {file}")
                continue

            obj = {"filename": file, "content": text}

            client.data_object.create(
                obj,
                class_name=CLASS_NAME,
                uuid=generate_uuid5(file + text)
            )

            print(f"📄 Ingested: {file}")

if __name__ == "__main__":
    ingest_docs()
    print("✅ All documents ingested successfully!")
